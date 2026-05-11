import os
import random
import io
import json
import anthropic
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from telegram import Update, ReplyKeyboardMarkup, InlineKeyboardMarkup, InlineKeyboardButton
from telegram.ext import ApplicationBuilder, MessageHandler, CommandHandler, CallbackQueryHandler, filters, ContextTypes

TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
WEATHER_API_KEY = os.getenv("WEATHER_API_KEY")
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")
UNSPLASH_API_KEY = os.getenv("UNSPLASH_API_KEY")

client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

user_histories = {}
user_notes = {}
user_quiz = {}
user_names = {}

CITY_MAP = {
    "toshkent": "Tashkent", "toshkend": "Tashkent",
    "samarqand": "Samarkand", "namangan": "Namangan",
    "andijon": "Andijan", "farg'ona": "Fergana", "fargona": "Fergana",
    "buxoro": "Bukhara", "nukus": "Nukus", "qarshi": "Qarshi",
    "termiz": "Termez", "jizzax": "Jizzakh", "navoiy": "Navoi",
    "urganch": "Urgench", "guliston": "Guliston",
}

MENU = ReplyKeyboardMarkup([
    ["🌤 Ob-havo", "💰 Valyuta kursi"],
    ["📝 Eslatmalar", "🗑 Eslatmani o'chirish"],
    ["🌐 Yangiliklar", "🎮 O'yin"],
    ["📊 Slayd yasash", "🤖 AI suhbat"],
    ["ℹ️ Yordam"]
], resize_keyboard=True)

QUIZ_QUESTIONS = [
    # Geografiya
    {"q": "O'zbekistonning poytaxti?", "options": ["Samarqand", "Toshkent", "Buxoro", "Namangan"], "answer": 1},
    {"q": "Dunyo eng baland tog'i?", "options": ["K2", "Kangchenjunga", "Everest", "Lhotse"], "answer": 2},
    {"q": "Eng katta okean?", "options": ["Atlantika", "Hind", "Tinch", "Arktika"], "answer": 2},
    {"q": "Qaysi mamlakat eng ko'p aholiga ega?", "options": ["Hindiston", "Xitoy", "AQSh", "Rossiya"], "answer": 0},
    {"q": "Afrika eng uzun daryosi?", "options": ["Kongo", "Niger", "Nil", "Zambezi"], "answer": 2},
    {"q": "Qaysi shahar Fransiyaning poytaxti?", "options": ["Berlin", "London", "Madrid", "Paris"], "answer": 3},
    {"q": "Dunyo eng katta mamlakatı?", "options": ["Kanada", "Xitoy", "AQSh", "Rossiya"], "answer": 3},
    {"q": "O'zbekistonda nechta viloyat bor?", "options": ["10", "12", "14", "16"], "answer": 2},
    {"q": "Qaysi tog' tizimi Yevropa va Osiyoni ajratadi?", "options": ["Alp", "Kavkaz", "Ural", "Karpat"], "answer": 2},
    {"q": "Yaponiyaning poytaxti?", "options": ["Osaka", "Kioto", "Tokyo", "Hiroshima"], "answer": 2},
    # Tarix
    {"q": "O'zbekiston mustaqillikni qachon oldi?", "options": ["1990", "1991", "1992", "1993"], "answer": 1},
    {"q": "Amir Temur qaysi yilda tug'ilgan?", "options": ["1320", "1336", "1350", "1362"], "answer": 1},
    {"q": "Birinchi jahon urushi qachon boshlangan?", "options": ["1912", "1913", "1914", "1915"], "answer": 2},
    {"q": "Ikkinchi jahon urushi qachon tugagan?", "options": ["1943", "1944", "1945", "1946"], "answer": 2},
    {"q": "Buyuk Ipak yo'li qaysi shahardan o'tgan?", "options": ["Buxoro", "Samarqand", "Xiva", "Hammasi"], "answer": 3},
    {"q": "Qaysi imperiya eng katta bo'lgan?", "options": ["Rim", "Britaniya", "Mo'g'ul", "Usmonli"], "answer": 1},
    {"q": "Kolumb Amerikani qachon kashf etgan?", "options": ["1488", "1490", "1492", "1495"], "answer": 2},
    # Fan va texnologiya
    {"q": "Quyosh sistemasida nechta sayyora?", "options": ["7", "8", "9", "10"], "answer": 1},
    {"q": "Internetni kim ixtiro qilgan?", "options": ["Bill Gates", "Steve Jobs", "Tim Berners-Lee", "Elon Musk"], "answer": 2},
    {"q": "Suvning kimyoviy formulasi?", "options": ["CO2", "H2O", "O2", "NaCl"], "answer": 1},
    {"q": "Eng tez hayvon?", "options": ["Sher", "Gepard", "Ot", "Burgut"], "answer": 1},
    {"q": "Inson tanasida nechta suyak bor?", "options": ["186", "196", "206", "216"], "answer": 2},
    {"q": "Qaysi planet Quyoshga eng yaqin?", "options": ["Venera", "Merkuriy", "Mars", "Yer"], "answer": 1},
    {"q": "Yorug'lik tezligi (km/s)?", "options": ["200,000", "250,000", "300,000", "350,000"], "answer": 2},
    {"q": "DNA nima?", "options": ["Oqsil", "Irsiy ma'lumot", "Vitamin", "Mineral"], "answer": 1},
    {"q": "Qaysi element eng yengil?", "options": ["Geliy", "Vodorod", "Litiy", "Azot"], "answer": 1},
    {"q": "Kompyuterni kim ixtiro qilgan?", "options": ["Bill Gates", "Alan Turing", "Steve Jobs", "Charles Babbage"], "answer": 3},
    # Sport
    {"q": "FIFA Jahon chempionati qancha yilda bir bo'ladi?", "options": ["2", "3", "4", "5"], "answer": 2},
    {"q": "Tennis kortida nechta o'yinchi o'ynaydi (yakka)?", "options": ["1", "2", "3", "4"], "answer": 1},
    {"q": "Olimpiya o'yinlari qancha yilda bir bo'ladi?", "options": ["2", "3", "4", "5"], "answer": 2},
    {"q": "Basketbol to'pi nechta oyoq balandlikdan tashlashda 3 ochko?", "options": ["4.5m", "5.8m", "6.75m", "7m"], "answer": 2},
    {"q": "Futbol darvozasining eni?", "options": ["6m", "7m", "7.32m", "8m"], "answer": 2},
    # Matematika
    {"q": "17 x 8 = ?", "options": ["126", "136", "146", "156"], "answer": 1},
    {"q": "√144 = ?", "options": ["11", "12", "13", "14"], "answer": 1},
    {"q": "2^10 = ?", "options": ["512", "1024", "2048", "256"], "answer": 1},
    {"q": "Uchburchak burchaklari yig'indisi?", "options": ["90°", "180°", "270°", "360°"], "answer": 1},
    {"q": "Pi soni taxminan?", "options": ["3.14", "3.16", "3.18", "3.12"], "answer": 0},
    # Adabiyot va madaniyat
    {"q": "Alisher Navoiy qaysi asarni yozgan?", "options": ["Shohnoma", "Xamsa", "Boburnoma", "Qutadg'u bilig"], "answer": 1},
    {"q": "Hamza Hakimzoda Niyoziy kim?", "options": ["Shoир", "Dramaturg", "Ikkalasi ham", "Rassом"], "answer": 2},
    {"q": "O'zbek tilida nechta harf bor (lotin)?", "options": ["26", "29", "32", "35"], "answer": 1},
    {"q": "Bobur mirzo qaysi davlatni tuzgan?", "options": ["Safaviy", "Boyqaro", "Boburiylar", "Shayboniy"], "answer": 2},
    # Texnologiya
    {"q": "Python qaysi yilda yaratilgan?", "options": ["1985", "1989", "1991", "1995"], "answer": 2},
    {"q": "Google qachon tashkil topgan?", "options": ["1996", "1997", "1998", "1999"], "answer": 2},
    {"q": "iPhone qachon chiqgan?", "options": ["2005", "2006", "2007", "2008"], "answer": 2},
    {"q": "WWW kim ixtiro qilgan?", "options": ["Bill Gates", "Tim Berners-Lee", "Vint Cerf", "Steve Jobs"], "answer": 1},
    {"q": "RAM nima?", "options": ["Doimiy xotira", "Operativ xotira", "Protsessor", "Disk"], "answer": 1},
]

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_names.pop(user_id, None)
    context.user_data["waiting_name"] = True
    greetings = [
        "Assalomu alaykum! Xush kelibsiz! 🎉\n\nMen sizning shaxsiy AI yordamchingizman! 🤖\n\nSizni to'g'ri chaqirish uchun — ismingizni ayting-chi? 😊",
        "Salom salom! Kelaveringlar, joy bor! 😄\n\nMen sizning AI yordamchingizman! 🤖\n\nQanday ism bilan murojaat qilay sizga? 🤔",
        "Voy, yangi mehmon! Xush kelibsiz! 🥳\n\nMen sizxizmatida turgan AI yordamchiman! 🤖\n\nAvval ismingizni biling-chi? 😊",
    ]
    import random as rnd
    await update.message.reply_text(rnd.choice(greetings))

async def get_weather(city: str) -> str:
    try:
        city_en = CITY_MAP.get(city.lower(), city)
        res = requests.get(f"http://api.openweathermap.org/data/2.5/weather?q={city_en}&appid={WEATHER_API_KEY}&units=metric").json()
        if res.get("cod") != 200:
            return f"❌ '{city}' topilmadi.\n\nSinab ko'ring: Toshkent, Namangan, Samarqand"
        return (
            f"🌤 *{res['name']}* ob-havosi:\n\n"
            f"🌡 Harorat: *{res['main']['temp']}°C* (his: {res['main']['feels_like']}°C)\n"
            f"💧 Namlik: *{res['main']['humidity']}%*\n"
            f"💨 Shamol: *{res['wind']['speed']} m/s*\n"
            f"📋 Holat: {res['weather'][0]['description']}"
        )
    except:
        return "❌ Ob-havo ma'lumotini olishda xato."

async def get_currency() -> str:
    try:
        res = requests.get("https://cbu.uz/oz/arkhiv-kursov-valyut/json/").json()
        wanted = ["USD", "EUR", "RUB", "GBP", "CNY", "KZT"]
        lines = ["💰 *Markaziy bank kurslari:*\n"]
        for item in res:
            if item["Ccy"] in wanted:
                lines.append(f"*{item['Ccy']}* — {float(item['Rate']):,.2f} so'm")
        lines.append(f"\n📅 {res[0]['Date']}")
        return "\n".join(lines)
    except:
        return "❌ Valyuta kursini olishda xato."

async def search_web(query: str) -> str:
    try:
        res = requests.post("https://api.tavily.com/search", json={
            "api_key": TAVILY_API_KEY, "query": query, "max_results": 3
        }).json()
        results = res.get("results", [])
        if not results:
            return "❌ Natija topilmadi."
        lines = [f"🌐 *'{query}' natijalari:*\n"]
        for i, r in enumerate(results[:3], 1):
            lines.append(f"{i}. *{r['title']}*\n{r['content'][:200]}...\n")
        return "\n".join(lines)
    except:
        return "❌ Qidirishda xato."

def get_unsplash_image(query: str) -> bytes | None:
    try:
        res = requests.get(
            f"https://api.unsplash.com/search/photos?query={query}&per_page=1&orientation=landscape",
            headers={"Authorization": f"Client-ID {UNSPLASH_API_KEY}"}
        ).json()
        photos = res.get("results", [])
        if not photos:
            return None
        img_url = photos[0]["urls"]["regular"]
        img_data = requests.get(img_url).content
        return img_data
    except:
        return None

def create_pptx(topic: str, slides_data: list) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    ACCENT = RGBColor(80, 180, 255)
    WHITE = RGBColor(255, 255, 255)
    BG_COLOR = RGBColor(12, 12, 22)

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Qora fon (rasm bo'lmasa)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        # Unsplash rasm
        search_query = slide_data.get("image_query", topic)
        img_bytes = get_unsplash_image(search_query)

        if img_bytes:
            try:
                img_stream = io.BytesIO(img_bytes)
                # O'ng tomonga rasm qo'yish
                pic = slide.shapes.add_picture(img_stream, Inches(7.0), Inches(0), Inches(6.33), Inches(7.5))
            except Exception as e:
                print(f"Rasm qo'shishda xato: {e}")

        # Chap tomonda gradient overlay
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(8.5), Inches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = BG_COLOR
        overlay.line.fill.background()

        # Accent chiziq
        line = slide.shapes.add_shape(1, Inches(0.3), Inches(1.8), Inches(6.5), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # Slayd raqami
        num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(1), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{i+1:02d}"
        r.font.size = Pt(14)
        r.font.color.rgb = ACCENT
        r.font.bold = True

        # Sarlavha
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.4), Inches(6.5), Inches(1.3))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = slide_data["title"]
        r2.font.size = Pt(28) if i == 0 else Pt(22)
        r2.font.bold = True
        r2.font.color.rgb = WHITE

        # Kontent
        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(2.0), Inches(6.5), Inches(5.0))
        tf3 = content_box.text_frame
        tf3.word_wrap = True

        for j, point in enumerate(slide_data["points"]):
            p3 = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
            p3.space_before = Pt(10)
            r3 = p3.add_run()
            r3.text = f"▸  {point}"
            r3.font.size = Pt(16)
            r3.font.color.rgb = WHITE

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

async def make_pptx(topic: str) -> tuple:
    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2000,
            system="""Siz prezentatsiya mutaxassisisiz. 6 ta slayd uchun JSON tayyorlang.
Faqat JSON qaytaring. Format:
[
  {"title": "Sarlavha", "points": ["Nuqta 1", "Nuqta 2", "Nuqta 3"], "image_query": "inglizcha rasm qidirish so'zi"},
  ...
]
O'zbek tilida yozing. image_query inglizcha bo'lsin (Unsplash uchun).""",
            messages=[{"role": "user", "content": f"Mavzu: {topic}"}]
        )
        text = response.content[0].text.strip().replace("```json", "").replace("```", "").strip()
        slides_data = json.loads(text)
        pptx_file = create_pptx(topic, slides_data)
        return pptx_file, None
    except Exception as e:
        return None, f"❌ Slayd yaratishda xato: {e}"

async def start_quiz(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    q = random.choice(QUIZ_QUESTIONS)
    if user_id not in user_quiz:
        user_quiz[user_id] = {"score": 0}
    user_quiz[user_id]["answer"] = q["answer"]
    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton(q["options"][0], callback_data="quiz_0"),
         InlineKeyboardButton(q["options"][1], callback_data="quiz_1")],
        [InlineKeyboardButton(q["options"][2], callback_data="quiz_2"),
         InlineKeyboardButton(q["options"][3], callback_data="quiz_3")],
    ])
    await update.message.reply_text(
        f"🎮 *Viktorina!*\n\n❓ {q['q']}\n\n🏆 Ball: {user_quiz[user_id]['score']}",
        parse_mode="Markdown", reply_markup=keyboard
    )

async def quiz_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    if user_id not in user_quiz:
        await query.edit_message_text("❌ /start bosing.")
        return
    chosen = int(query.data.split("_")[1])
    if chosen == user_quiz[user_id]["answer"]:
        user_quiz[user_id]["score"] += 1
        result = f"✅ *To'g'ri!* +1 ball\n🏆 Jami: {user_quiz[user_id]['score']}"
    else:
        result = f"❌ *Noto'g'ri!*\n🏆 Jami: {user_quiz[user_id]['score']}"
    keyboard = InlineKeyboardMarkup([[InlineKeyboardButton("🔄 Keyingi", callback_data="quiz_next")]])
    await query.edit_message_text(result, parse_mode="Markdown", reply_markup=keyboard)

async def quiz_next_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    q = random.choice(QUIZ_QUESTIONS)
    if user_id not in user_quiz:
        user_quiz[user_id] = {"score": 0}
    user_quiz[user_id]["answer"] = q["answer"]
    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton(q["options"][0], callback_data="quiz_0"),
         InlineKeyboardButton(q["options"][1], callback_data="quiz_1")],
        [InlineKeyboardButton(q["options"][2], callback_data="quiz_2"),
         InlineKeyboardButton(q["options"][3], callback_data="quiz_3")],
    ])
    await query.edit_message_text(
        f"🎮 *Viktorina!*\n\n❓ {q['q']}\n\n🏆 Ball: {user_quiz[user_id]['score']}",
        parse_mode="Markdown", reply_markup=keyboard
    )

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text = update.message.text.strip()

    # Ism kutilayotgan bo'lsa
    if context.user_data.get("waiting_name"):
        user_names[user_id] = text
        context.user_data["waiting_name"] = False
        funny_welcomes = [
            f"Vay, {text} aka, zo'r ism ekan! 🔥\n\nEndi men sizningman! 😄 Nima qilishimni xohlaysiz?",
            f"Oho, {text}! Juda chiroyli ism! ✨\n\nKeling endi ishga kirishamiz! 💪",
            f"Xush kelibsiz, {text} aka! 🎉\n\nMen 24/7 xizmatingizdaman, hatto tunda ham! 😄",
            f"Zo'r, {text}! Endi biz do'stmiz! 🤝\n\nNimada yordam kerak?",
        ]
        import random as rnd
        await update.message.reply_text(
            rnd.choice(funny_welcomes),
            reply_markup=MENU
        )
        return

    # Ismni olish
    user_name = user_names.get(user_id, "Aka")

    if text == "🌤 Ob-havo":
        await update.message.reply_text("Qaysi shahar?\n\nMasalan: *Toshkent, Namangan*", parse_mode="Markdown")
        context.user_data["mode"] = "weather"
        return
    if text == "💰 Valyuta kursi":
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
        await update.message.reply_text(await get_currency(), parse_mode="Markdown", reply_markup=MENU)
        return
    if text == "🌐 Yangiliklar":
        await update.message.reply_text("Nima haqida qidirmoqchisiz?")
        context.user_data["mode"] = "search"
        return
    if text == "🎮 O'yin":
        await start_quiz(update, context)
        return
    if text == "📊 Slayd yasash":
        await update.message.reply_text("Qaysi mavzuda slayd?\n\nMasalan: *Sun'iy intellekt, Iqlim o'zgarishi*", parse_mode="Markdown")
        context.user_data["mode"] = "slides"
        return
    if text == "📝 Eslatmalar":
        notes = user_notes.get(user_id, [])
        if not notes:
            await update.message.reply_text("📝 Eslatma yo'q.\n\nQo'shish: *eslatma: matn*", parse_mode="Markdown")
        else:
            await update.message.reply_text("📝 *Eslatmalar:*\n\n" + "\n".join([f"{i+1}. {n}" for i, n in enumerate(notes)]), parse_mode="Markdown")
        return
    if text == "🗑 Eslatmani o'chirish":
        notes = user_notes.get(user_id, [])
        if not notes:
            await update.message.reply_text("Eslatma yo'q.")
        else:
            await update.message.reply_text("*Qaysi raqam?*\n\n" + "\n".join([f"{i+1}. {n}" for i, n in enumerate(notes)]), parse_mode="Markdown")
            context.user_data["mode"] = "delete_note"
        return
    if text == "🤖 AI suhbat":
        await update.message.reply_text("Savolingizni yozing! 🤖")
        context.user_data["mode"] = "ai"
        return
    if text == "ℹ️ Yordam":
        await update.message.reply_text(
            "ℹ️ *Yordam:*\n\n🌤 Ob-havo • 💰 Valyuta\n🌐 Yangiliklar • 🎮 Viktorina\n📊 PPTX slayd (rasimli!) • 📝 Eslatmalar\n🤖 AI suhbat\n\nYaratuvchi: @Samik_1806\nEslatma: *eslatma: matn*",
            parse_mode="Markdown", reply_markup=MENU
        )
        return
    if text.lower().startswith("eslatma:"):
        note = text[8:].strip()
        if note:
            if user_id not in user_notes:
                user_notes[user_id] = []
            user_notes[user_id].append(note)
            await update.message.reply_text(f"✅ Saqlandi: *{note}*", parse_mode="Markdown")
        return

    mode = context.user_data.get("mode", "ai")

    if mode == "weather":
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
        await update.message.reply_text(await get_weather(text), parse_mode="Markdown", reply_markup=MENU)
        context.user_data["mode"] = "ai"
        return
    if mode == "search":
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
        await update.message.reply_text(await search_web(text), parse_mode="Markdown", reply_markup=MENU)
        context.user_data["mode"] = "ai"
        return
    if mode == "slides":
        await update.message.reply_text("⏳ Slayd va rasmlar tayyorlanmoqda, biroz kuting...")
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_document")
        pptx_file, error = await make_pptx(text)
        if error:
            await update.message.reply_text(error, reply_markup=MENU)
        else:
            await update.message.reply_document(
                document=pptx_file,
                filename=f"{text[:30]}.pptx",
                caption=f"📊 *{text}* bo'yicha slayd tayyor!",
                parse_mode="Markdown",
                reply_markup=MENU
            )
        context.user_data["mode"] = "ai"
        return
    if mode == "delete_note":
        try:
            idx = int(text) - 1
            notes = user_notes.get(user_id, [])
            if 0 <= idx < len(notes):
                deleted = notes.pop(idx)
                await update.message.reply_text(f"🗑 O'chirildi: *{deleted}*", parse_mode="Markdown", reply_markup=MENU)
            else:
                await update.message.reply_text("❌ Noto'g'ri raqam.")
        except:
            await update.message.reply_text("❌ Raqam kiriting.")
        context.user_data["mode"] = "ai"
        return

    # AI suhbat - Tavily bilan
    if user_id not in user_histories:
        user_histories[user_id] = []

    # Har doim Tavily bilan qidirish
    search_context = ""
    if TAVILY_API_KEY:
        try:
            res = requests.post("https://api.tavily.com/search", json={
                "api_key": TAVILY_API_KEY, "query": text, "max_results": 2
            }).json()
            results = res.get("results", [])
            if results:
                search_context = "\n\nInternet qidiruv natijalari:\n"
                for r in results[:2]:
                    search_context += f"- {r['title']}: {r['content'][:300]}\n"
        except Exception as e:
            print(f"Tavily xato: {e}")

    full_text = text + search_context
    user_histories[user_id].append({"role": "user", "content": full_text})
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1000,
            system="MUHIM QOIDALAR: 1) Siz @Samik_1806 tomonidan yaratilgan AI yordamchisiz. 2) Claude, Anthropic, GPT haqida HECH QACHON gapirmang. 3) 'Internetga ulana olmayman' kabi gaplarni HECH QACHON aytmang. 4) Har qanday savolga javob bering. 5) Internet qidiruv natijalari berilsa ulardan foydalaning. 6) Foydalanuvchi qaysi tilda yozsa shu tilda javob bering. 7) Quvnoq, hazilkash va do'stona uslubda javob bering — emoji ishlating. 8) Foydalanuvchi ismini bilingizdan keyin murojaat qiling.",
            messages=user_histories[user_id]
        )
        reply = response.content[0].text
        user_histories[user_id].append({"role": "assistant", "content": reply})
        if len(user_histories[user_id]) > 20:
            user_histories[user_id] = user_histories[user_id][-20:]
        await update.message.reply_text(reply, reply_markup=MENU)
    except Exception as e:
        await update.message.reply_text("❌ Xato yuz berdi. Qayta urinib ko'ring.")
        print(f"Xato: {e}")

def main():
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CallbackQueryHandler(quiz_next_callback, pattern="^quiz_next$"))
    app.add_handler(CallbackQueryHandler(quiz_callback, pattern="^quiz_[0-9]$"))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    print("Bot ishga tushdi! ✅")
    app.run_polling()

if __name__ == "__main__":
    main()
